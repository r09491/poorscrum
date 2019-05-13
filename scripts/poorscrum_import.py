#!/usr/bin/env python
# -*- coding: utf-8 -*-

__author__ = "sepp.heid@t-online.de"
__doc__ = """ """

from pptx import Presentation

from poorscrum import Poorscrum_Tasks, Status, read_fields, EMPTY_PPTX, story_points, __version__

import configparser

import argparse
import os
import sys

import logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(os.path.basename(sys.argv[0]))

          
def parse_arguments():
    """Parse command line arguments"""

    parser = argparse.ArgumentParser(
        prog = os.path.basename(sys.argv[0]),
        description ="Imports the story info into the slides of  PPTX file from text files",
        epilog = "Ensure the '.pptx/pptx_pickup.ini' file exists in the home directory")

    parser.add_argument("--version", action="version", version=__version__)

    parser.add_argument("--dry", required=False,
                        action="store_true", 
                        help="Do not save the presentation")

    parser.add_argument("--empty", required=False,
                        action="store_true", 
                        help="Start with empty slide flag")

    parser.add_argument('--status_first', type=Status, choices=list(Status), default = "none",
                        help="Lowest story status considered")

    parser.add_argument('--status_last', type=Status, choices=list(Status), default = "out",
                        help="Highest story status considered")

    parser.add_argument("to_pptx", nargs=1, 
                        help="PPTX file to import the story slides to")

    parser.add_argument("from_text", nargs="+", 
                        help="Text files with story content")

    return parser.parse_args()


def import_text(to_frame, text):
    p = to_frame.paragraphs[0] # Already there!
    r = p.add_run()

    in_hlink = False
    hlink = ""
    old = ""
    for new in text:
        if in_hlink: ## in hyperlink text
            if new == '>' :
                in_hlink = False                

                if old == "/":
                    r = p.add_run()
                else:
                    r = p.add_run()
                    r.hyperlink.address = hlink

            elif new == '"': 
                pass         ## Skip: Powerpoint does not like it 

            elif new == '/': 
                hlink += '/' ## Needed for directory paths
            else:            
                hlink += new ## Add anything else to the address part
                
        else: ## standard text
            if new == '<':
                ## A run already exists!
                in_hlink = True                
                hlink = ""

            elif new == "\n":
                p = to_frame.add_paragraph()
                r = p.add_run()

            else:
                r.text += new ## Add anything else to the text part

        old = old if new.isspace() else new  ## Transitions without space
        

def import_story(from_story, to_prs, with_pickup):
    """ Imports the story items for a story slide
    """

    """!!! Our layout is always the last in the layout list !!!"""
    story_slide_layout = to_prs.slide_layouts[-1]
    to_slide = to_prs.slides.add_slide(story_slide_layout)

    for item, index in with_pickup.items():
        try:
            shape = to_slide.placeholders[int(index)]
        except KeyError:
            return None
        
        if not shape.has_text_frame:
            continue

        try:
            item_text = from_story.get(item, "text") 
        except configparser.NoSectionError:
            return None

        import_text(shape.text_frame, item_text)

    return to_slide


def import_tasks(from_tasks, to_prs):
    """ Imports the tasks items for a tasks slide
    """

    """!!! Our tasks layout is always the one before the story layout !!!"""
    tasks_slide_layout = to_prs.slide_layouts[-2]
    to_slide = to_prs.slides.add_slide(tasks_slide_layout)
    to_table = Poorscrum_Tasks(to_slide, to_prs.slide_height, to_prs.slide_width)
    lastrow = len(to_table.table.rows)-1
    
    try:
        tasks = from_tasks.items("tasks")
    except:
        return None

    wstart, wleft, wdone = 0, 0, 0
    for row, task in enumerate(tasks):
        if row < lastrow:
            # Calculate all rows
            data = task[1].split(',')
            wstart += int(data[1]) 
            wleft += int(data[2])
            wdone += int(data[3])
            # Output leading tasks only
            to_table.put_as_row(row, data)

    # Fibonacci!
    total = ["Total",
             str(story_points(wstart)),
             str(story_points(wleft)),
             str(story_points(wdone)), "Points"]
    to_table.put_as_row(lastrow, total)
        
    return to_slide


def main():
    args = parse_arguments()

    if args.to_pptx is None:
        logger.error("Must provide target PPTX files for slides!")
        return 1

    if os.path.splitext(args.to_pptx[0])[1] != ".pptx":
        logger.error("Must provide '.pptx' extension for presentation!")
        return 2

    if args.from_text is None:
        logger.error("Must provide TEXT backlog files for input!")
        return 3

    fields_map = read_fields()
    if fields_map is None:
        logger.error("Must provide correct story item mappings!")
        return 4

    input_pptx = EMPTY_PPTX if args.empty else args.to_pptx[0]

    try:
        to_prs = Presentation(input_pptx)
    except:
        logger.error("Presentation '{}' is not defined".format(input_pptx))
        return 5

    logger.info("PPTX Backlog has '{:d}' slides before import"
                .format(len(to_prs.slides)))


    for story_file in args.from_text:
        name, ext = os.path.splitext(os.path.basename(story_file))

        if ext != ".story":
            logger.error("Must provide story file names with '.story' extension!")
            return 6

        try:
            num = int(name.split('_')[0])
        except:
            logger.error("Must provide story file names with leading digits!")
            return 7


        """ Read the story """
        
        story = configparser.ConfigParser()
        story.read(story_file)

        """ abort if the story is not in desired state """

        try:
            slide_status = Status(story.get("status", "text"))
        except:
            logger.warn("Skipped the file prefixed #{:d}. Wrong Status Item!"
                        .format(num))
            return 8

        """ Skip the story is not in desired range """
        
        first_status = Status(args.status_first)
        last_status = Status(args.status_last)
        if  slide_status < first_status:
            logger.info("Skipped the slide #{:d}. Status '{}' is below '{}'."
                        .format(num, slide_status, first_status))
            continue
        if  slide_status > last_status:
            logger.info("Skipped the slide #{:d}. Status '{}' is after '{}'."
                        .format(num, slide_status, last_status))
            continue

        """ Append the story to the presentation """
        
        slide = import_story(story, to_prs, fields_map)
        if slide is None:
            logger.error("Wrong text story format: '{}'".format(story_file))
            return 9

        tasks = import_tasks(story, to_prs)
        if tasks is None:
            logger.warn("There are no tasks in: '{}'".format(story_file))

        logger.info("Import ok: '{}'".format(story_file))


    if not args.dry:
        """ Save the presentation """

        try:
            to_prs.save(args.to_pptx[0]) 
        except:
            logger.error("Cannot save the presentation. If it is open consider to close it!")
            return 10

        logger.info("Presentation saved to '{}' .".format(args.to_pptx[0]))

    else:
        
        logger.info("Would save presentation to '{}' .".format(args.to_pptx[0]))

    return 0

        
if __name__ == "__main__":
    sys.exit(main())

