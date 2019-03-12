# -*- coding: utf-8 -*-

from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR

class Poorscrum_Tasks:


    def __init__(self, slide, slide_height, slide_width, 
                 top = Cm(0.5), left = Cm(0.5), rows = 8, cols = 5):

        height = int(slide_height - 2*top)
        width = int(slide_width/2 - 2*left)
        width_task = int(0.80 * width)
        width_work = int(0.05 * width)
        width_devs = int(0.15 * width)
        task_table = slide.shapes.add_table(rows, cols, left, top, width, height)

        """ No special treatment of the first and last rows and coloumns """
        
        table = task_table.table
        table.first_row = False
        table.first_col = False
        table.last_row = False
        table.last_col = False

        """ Set the width of the columns """
        
        table.columns[0].width = int(0.55*width)
        table.columns[1].width = int(0.1*width)
        table.columns[2].width = int(0.1*width)
        table.columns[3].width = int(0.1*width)
        table.columns[4].width = int(0.15*width)

        default_text = ["<task{:d}>", "#e", "#l", "#d", "<dev>"]
        
        for rkey, row in enumerate(table.rows):
            
            for ckey, (cell, text) in enumerate(zip(row.cells, default_text)):
                t = cell.text_frame
                if ckey == 0:
                    t.word_wrap = True
                    t.auto_size = MSO_AUTO_SIZE.NONE
                    t.vertical_anchor = MSO_ANCHOR.TOP
                else:
                    t.word_wrap = False
                    t.auto_size = MSO_AUTO_SIZE.NONE
                    t.vertical_anchor = MSO_ANCHOR.MIDDLE

                p = t.paragraphs[0]
                if ckey == 0:
                    p.alignment = PP_ALIGN.LEFT
                else:
                    p.alignment = PP_ALIGN.CENTER

                r = p.add_run()
                r.text = default_text[ckey].format(rkey)
                r.font.size = Pt(10)
                r.font.name = 'Consolas'

        self.table = table

        
    def put_as_row(self, row, task):

        for cell, text in zip(self.table.rows[row].cells, task):
            t = cell.text_frame
            p = t.paragraphs[0]
            r = p.runs[0]
            r.text = text

            
    def put_as_list(self, tasks):

        for row, task in zip(self.table.rows, tasks):
            self.put_as_row(row, task)
            

    def get_from_row(self, row):
        return None

    
    def get_from_list(self):
        return None
    
