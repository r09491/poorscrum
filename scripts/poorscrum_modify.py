#!/usr/bin/env python
# -*- coding: utf-8 -*-

__author__ = "sepp.heid@t-online.de"
__doc__ = """ """


from pptx import Presentation
from pptx.exc import PackageNotFoundError

from poorscrum import Status, read_fields, __version__

import configparser

import argparse

import os
import sys

import logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(os.path.basename(sys.argv[0]))


def parse_arguments():
    """Parse command line arguments"""

    parser = argparse.ArgumentParser(
        prog = os.path.basename(sys.argv[0]),
        description = "Modify a single exported story from a slide of PPTX file",
        epilog = "Ensure the '.pptx/pptx_pickup.ini' file exists in the the home directory")
    

    parser.add_argument("--version", action = "version", version=__version__)

    parser.add_argument("--dry", required=False,
                        action="store_true", 
                        help="Do not create text stories")

    parser.add_argument('--field', required=False, default=None,
                        help="The only field for update")

    parser.add_argument("from_pptx", nargs="?", 
                        help="The name of the PPTX file with story slides to modify")

    parser.add_argument("to_story_file", nargs="?", 
                        help="The target story file. Must exist!")

    return parser.parse_args()


def export_story_as_config(from_slide, with_fields):
    """ Returns the story items for a story slide in the config format

    Each slide contains placeholder with a text frame which in turn contains
    paragraphs with text. Only ASCII characters are allowed in the story text
    files to overcome the different encodings on computers.

    Each placeholder is associated with a story item. They become a 'section' in
    the config file with the 'key' text associated with the parsed text.
    """
    
    story = configparser.RawConfigParser()

    for item, index in with_fields.items():
        try:
            shape = from_slide.placeholders[int(index)]
        except KeyError:
            return None
        
        if not shape.has_text_frame:
            continue

        text = ""
        paras = shape.text_frame.paragraphs
        for p in paras:
            for r in p.runs:
                link = r.hyperlink
                has_link = link is not None and (link.address is not None)
                if has_link:
                    text += '<'
                    for c in link.address:
                        text += c if ord(c) < 128 else ''
                    text += '>'
                for c in r.text:
                    text += c if ord(c) < 128 else ''
                if has_link:
                    text += '</>'
            if( len(paras) > 1):
                text += "\n"
                    
        story.add_section(item)
        story.set(item, "text", text.strip())

    return story


def main():
    args = parse_arguments()

    if args.from_pptx is None:
        logger.error("Must provide PPTX backlog file name for input!")
        return 1

    if os.path.splitext(args.from_pptx)[1] != ".pptx":
        logger.error("Must provide '.pptx' extension for presentation!")
        return 2

    if (args.to_story_file is None):
        logger.error("Must provide target file for story!")
        return 3

    if not os.path.isfile(args.to_story_file):
        logger.error("Story file '{}' does not exist!".format(args.to_story_file))
        return 4

    logger.info("Story file '{}' found!".format(args.to_story_file))

    try:
        slidenum = int(int(os.path.basename(args.to_story_file).split('_')[0])/10)
    except:
        logger.error("Cannot retrieve slide number!")
        return 5
        
    fields_map = read_fields()
    if fields_map is None:
        logger.error("Must provide correct story field mappings. Run 'pptx_learn.py'!")
        return 6

    try:
        prs = Presentation(args.from_pptx)
    except NameError:
        logger.error("Presentation has illegal name.")
        return 7
    except PackageNotFoundError:
        logger.error("Presentation is not found.")
        return 8
    
    logger.info("The slide number used for modification is '{:d}'".format(slidenum))

    if slidenum >=  len(prs.slides):
        logger.error("The slide number {:d} is illegal.".format(slidenum))
        return 9

    if args.field is not None:
        if not (args.field in fields_map.keys()):
            logger.error("The field '{}' is illegal.".format(args.field))
            return 10

        logger.info("About to modify the field '{}'.".format(args.field))
    else:
        logger.info("All fields of slide '{:d}' are taken into account.".format(slidenum))


    """ Read the story from the slide """        

    story_slide = None
    for num, slide in enumerate(prs.slides):
        if num == slidenum-1:        
            story_slide = export_story_as_config(slide, fields_map)
            break
    if story_slide is None:
        logger.info("Skipped the slide #{:d} with foreign format.".format(num))
        return 11

    if args.field is None:
        """ Update the complete story """
        
        if not args.dry:
            with open(args.to_story_file, 'w') as storyfile:
                story_slide.write(storyfile)
            logger.info("Saved the story of slide #{:d} as '{}'."
                        .format(num, args.to_story_file))
        else:
            logger.info("Would save the story of slide #{:d} as '{}'."
                        .format(num, args.to_story_file))

    else:
        """ Update the field in the story"""        

        story_slide_text = story_slide.get(args.field, "text")

        story_config = configparser.ConfigParser()
        story_config.read(args.to_story_file)
        story_config_text = story_config.get(args.field, "text")
        story_config.set(args.field, "text", story_slide_text)

        logger.info("'{}' replaces '{}'."
                    .format(story_config_text, story_slide_text))

        if not args.dry:
            with open(args.to_story_file, 'w') as storyfile:
                story_config.write(storyfile)
                logger.info("Modified the story of slide #{:d} as '{}'."
                            .format(num, args.to_story_file))
        else:
            logger.info("Would modify the story of slide #{:d} as '{}'."
                        .format(num+1, args.to_story_file))
        
    return 0

        
if __name__ == "__main__":
    sys.exit(main())

