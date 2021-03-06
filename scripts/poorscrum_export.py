#!/usr/bin/env python
# -*- coding: utf-8 -*-

__author__ = "sepp.heid@t-online.de"
__doc__ = """ """


from pptx import Presentation
from pptx.exc import PackageNotFoundError

from poorscrum import Status, read_fields, story_points, __version__ 

import configparser

import argparse

import os
import sys

import logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(os.path.basename(sys.argv[0]))


def parse_arguments():
    """Parse command line arguments"""

    parser = argparse.ArgumentParser(
        prog = os.path.basename(sys.argv[0]),
        description = "Export stories from text file to the slides of PPTX file",
        epilog = "Ensure the '.pptx/pptx_pickup.ini' file exists in the the home directory")
    

    parser.add_argument("--version", action = "version", version=__version__)

    parser.add_argument("--dry", required=False,
                        action="store_true", 
                        help="Do not create text stories")

    parser.add_argument("from_pptx", nargs="?", 
                        help="The name of the PPTX file with story slides to export")

    parser.add_argument("to_slide_dir", nargs="?", 
                        help="The target directory for the text story files. Must not exist!")


    kgroup = parser.add_argument_group(title="export to KANBAN")
    kgroup.add_argument("--kanban", required=False,
                        action="store_true", 
                        help="Export to KANBAN directoy structure")

    fgroup = parser.add_argument_group(title="filter stories")
    fgroup.add_argument('--status_first', type=Status, choices=list(Status), default = "none",
                        help="Lowest story status to be considered")

    fgroup.add_argument('--status_last', type=Status, choices=list(Status), default = "out",
                        help="Highest story status to be considered")


    ngroup = parser.add_argument_group(title="name stories")
    ngroup.add_argument("--with_title", type=int, default=35,
                        help="Add a part of the TITLE field to the story file name")
    ngroup.add_argument("--with_values", required=False,
                        action="store_true", 
                        help="Prefix the story file name with the VALUE field for priority")
    ngroup.add_argument("--with_ids", required=False,
                        action="store_true", 
                        help="Prefix the story file name with the ID field")
    ngroup.add_argument("--skip_counts", required=False,
                        action="store_true", 
                        help="Skip the count prefix for the file name")

    return parser.parse_args()


def extract_text(from_frame):
    text = ""
    paras = from_frame.paragraphs
    for p in paras:
        for r in p.runs:
            link = r.hyperlink
            has_link = link is not None and (link.address is not None)
            if has_link:
                text += '<'
                for c in link.address:
                    text += c if ord(c) < 128 else ''
                text += '>'
            for c in r.text:
                text += c if ord(c) < 128 else ''
            if has_link:
                text += '</>'
        if( len(paras) > 1):
            text += "\n"
    return text


def export_story_to_config(from_slide, with_fields):
    """ Returns the story items for a story slide in the config format

    Each story slide contains placeholders with a text frame which in turn
    contains paragraphs with text. Only ASCII characters are allowed in the
    story text files to overcome the different encodings on computers.

    Each placeholder is associated with a story item. They become a 'section'
    in the config file with the 'key' text associated with the parsed text.
    """
    
    story = configparser.RawConfigParser()

    for item, index in with_fields.items():
        try:
            shape = from_slide.placeholders[int(index)]
        except KeyError:
            return None

        if not shape.has_text_frame:
            continue

        text = extract_text(shape.text_frame)
                    
        story.add_section(item)
        story.set(item, "text", text.strip())

    return story


def append_tasks_to_config(from_slide, to_story):
    
    """ Find the container of the table.
    Abort if there is more than one!  """
    for shape in from_slide.shapes:
        if not shape.has_table:
            return None

    table = shape.table
    lastrow = len(table.rows)-1

    """ Add each task in the slide in tasks section """
    wstart, wleft, wdone = 0, 0, 0
    to_story.add_section("tasks")
    for num, row in enumerate(table.rows):
        if num < lastrow:
            tasks = [extract_text(cell.text_frame) for cell in row.cells]
            to_story.set("tasks","task{:d}".format(num+1),",".join(tasks))
            wstart += int(tasks[1]) 
            wleft += int(tasks[2])
            wdone += int(tasks[3])
    total = ",".join(["Total",
                      str(story_points(wstart)),
                      str(story_points(wleft)),
                      str(story_points(wdone)), "Points"])
    to_story.set("tasks", "total", total)
    return to_story


def write_config(num, count, story, args):
    """ Should have been done by argparse !"""
    first_status = Status(args.status_first)
    last_status = Status(args.status_last)

    try:
        slide_status = Status(story.get("status", "text"))
    except:
        """ Any illegal status is considered as out """
        slide_status = last_status

    if  slide_status < first_status:
        logger.info("Skipped the slide #{:d}. Status '{}' is below '{}'."
                    .format(num, slide_status, first_status))
        return None
        
    if  slide_status > last_status:
        logger.info("Skipped the slide #{:d}. Status '{}' is after '{}'."
                    .format(num, slide_status, last_status))
        return None

    """ Store in a state dependent directory structure (KANBAN) """

    to_slide_dir = args.to_slide_dir
    if args.kanban:
        to_slide_dir = os.path.join(to_slide_dir, slide_status.value.lower())

        if (slide_status == Status("ANALYSING")) or (slide_status == Status("SPRINTING")):
            """ A Dev owns a story only during  ANALYSING and SPRINTING. In the
            other states a story belongs to the teams
            """
                
            """ Add the first developper if available """
            slide_dev = story.get("devs", "text").split()
            try:
                to_slide_dir = os.path.join(to_slide_dir, slide_dev[0])
            except:
                logger.error("The slide #{:d} has no Dev. Aborted!".format(num))
                return None

    if not os.path.exists(to_slide_dir):
        try:
            os.makedirs(to_slide_dir)
        except:
            logger.error("Failed to create the kanban directory '{}'.".
                         format(to_slide_dir))
        logger.info("Created the kanban directory '{}'.".
                    format(to_slide_dir))


    """ Generate a very special file name from title to guess content"""
        
    if args.with_title == 0:
        """ Provide a pure numeric file name """
        story_filename = "{:04d}".format(10*count)
    else:
        """ Provide a narrative file name """
        title = story.get("title", "text").strip()
        title = title if len(title) < args.with_title else title[:args.with_title]
        story_filename = title.strip().lower().replace(' ', '_').replace('/', '_')

        """ Prefix by the count if not explicitly skipped """
        if not args.skip_counts:
            story_filename = "{:04d}_{}".format(10*count, story_filename)

    """ Append the extension """
    story_filename = "{}.story".format(story_filename)

    if args.with_ids:
        try:
            id = int(story.get("id", "text"))
        except:
            """ Any illegal id becomes no id """
            id = 0
        story_filename = "{:04d}_{}".format(id, story_filename)

    if args.with_values:
        try:
            value = int(story.get("value", "text"))
        except:
            """ Any illegal value becomes no value """
            value = 0
        priority = 0 if value < 0 or value > 100 else 100-value
        story_filename = "{:03d}_{}".format(priority, story_filename)


    """ Save the presentation """
        
    story_pathname = os.path.join(to_slide_dir, story_filename) 

    if not args.dry:
        with open(story_pathname, 'w') as storyfile:
            story.write(storyfile)
        logger.info("Saved the story slide #{:d} as '{}'."
                    .format(num, story_pathname))
    else:
        logger.info("Would save the story of slide #{:d} as '{}'."
                    .format(num, story_pathname))
            
    return story_pathname


def main():
    args = parse_arguments()

    if args.from_pptx is None:
        logger.error("Must provide PPTX backlog file name for input!")
        return 1

    if os.path.splitext(args.from_pptx)[1] != ".pptx":
        logger.error("Must provide '.pptx' extension for presentation!")
        return 2

    if (args.to_slide_dir is None):
        logger.error("Must provide target directory for slides!")
        return 3

    if ('*' in  args.to_slide_dir) or \
           ('%' in  args.to_slide_dir) :
        logger.error("Must provide target directory without wild cards!")
        return 4

    if args.with_values and args.kanban:
        logger.error("Must provide legal parameter combinations!")
        return 5
        
    fields_map = read_fields()
    if fields_map is None:
        logger.error("Must provide correct story field mappings. Run 'pptx_learn.py'!")
        return 6


    try:
        prs = Presentation(args.from_pptx)
    except NameError:
        logger.error("Presentation has illegal name.")
        return 7
    except PackageNotFoundError:
        logger.error("Presentation is not found.")
        return 8
    
    logger.info("The Backlog '{}' has '{:d}' slide(s)."
                .format(os.path.basename(args.from_pptx), len(prs.slides)))


    if not args.dry:
        try:
            os.makedirs(args.to_slide_dir)
        except:
            logger.error("Failed to create the slide directory '{}'. Consider to delete it!"
                         .format(args.to_slide_dir))
            return 9

        logger.info("Created the directory '{}' for the story files"
                    .format(args.to_slide_dir))
    else:
        logger.info("Would create the directory '{}' for the story files"
                    .format(args.to_slide_dir))

    story = None    
    selected_stories = 0
    for num, slide in enumerate(prs.slides):

        """
        A story consists of the a slide pair with the specification fields and a
        slide with the implementation tasks.

        First the specification is exported to the config file. If ok then the
        tasks are appended. So for the two slides there is one config file.
        """
        if story:
            """ Try to add tasks. No problem if there is none"""
            tasks = append_tasks_to_config(slide, story)
            if tasks is None:
                logger.warn("Tasks slide is missing after story slide #{:d}.".format(num))


            """ Always save it """
            if not args.dry:
                write_config(num, selected_stories, story, args)

            if tasks:
                logger.info("Task slide #{:d} is appended to file.".format(num+1))
                story = None
                continue
            
        story = export_story_to_config(slide, fields_map)
        if story is None:
            logger.info("Skipped the slide #{:d} as story slide.".format(num+1))
            continue

        selected_stories += 1

    if not args.dry:
        if story:
            tasks = append_tasks_to_config(slide, story)
            if tasks is None:
                logger.warn("Tasks slide is missing after story slide #{:d}.".format(num+1))

            write_config(num, selected_stories, story, args)
            
        logger.info("Saved '{:d}' stories with tasks.".format(selected_stories))
    else:
        logger.info("Would have saved {:d} stories.".format(selected_stories))
        
    return 0

        
if __name__ == "__main__":
    sys.exit(main())

