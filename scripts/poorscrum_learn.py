#!/usr/bin/env python
# -*- coding: utf-8 -*-

__author__ = "sepp.heid@t-online.de"
__doc__ = """ """

import argparse
import os
import sys

from pptx import Presentation

import poorscrum

import logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(os.path.basename(sys.argv[0]))


def parse_arguments():
    """Parse command line arguments"""

    parser = argparse.ArgumentParser(
        prog=os.path.basename(sys.argv[0]),
        description="Extracts the item indices of PPTX file",
        epilog = """
        Extracts the item indices of PPTX file

        To be run if the master of the PPTX file changed:
        (1) Insert new slide into PPTX with the desired layout,
        (2) Enter a serial number in the fields of the slide starting from 0,
        (3) Setup the INI file with the placeholder name list
        (4) Run the script
        """)
    
    parser.add_argument('--version', action='version', version=poorscrum.__version__)

    parser.add_argument('from_pptx', nargs='?', 
                        help='PPTX file with story slides with the destinct layput')
    parser.add_argument('and_slide', nargs='?', type=int, 
                        help='Slide # to be used for indexing with default content')

    return parser.parse_args()


def extract_story_as_field_map(from_slide, use_pickup):
    """ Extracts the story items from a story alide
    """

    keys = list(use_pickup.keys())
    for shape in from_slide.placeholders:
        if not shape.has_text_frame:
            continue
        """ Get the number in the place holder """
        text = ""
        for p in shape.text_frame.paragraphs:
            for c in p.text:
                text += c if ord(c) < 128 else ''
        try:
            use_pickup[keys[int(text)]] = \
                str(shape.placeholder_format.idx)
        except:
            return None

    return use_pickup


def main():
    args = parse_arguments()

    if args.from_pptx is None:
        logger.error("Must provide PPTX backlog file name for input")
        return 1

    if os.path.splitext(args.from_pptx)[1] != ".pptx":
        logger.error("Must provide '.pptx' extension for presentation!")
        return 2

    if args.and_slide is None:
        logger.error("Must provide slide number as index output")
        return 3

    fields_map = poorscrum.read_fields()
    if fields_map is None:
        logger.error("Must provide initial pickup file")
        return 4
        
    try:
        prs = Presentation(args.from_pptx)
    except NameError:
        logger.error("Presentation is not defined")
        return 5

    if not (args.and_slide in range(1,len(prs.slides)+1)):
        logger.error("Must provide slide number within range of PPTX (1 to {:d})"
                     .format(len(prs.slides)))
        return 6

    fields_map = extract_story_as_field_map(prs.slides[args.and_slide-1], fields_map)
    if fields_map is None:
        logger.error("Must provide correct initial pickup file")
        return 7

    pickup_file = poorscrum.write_fields(fields_map)

    logger.info("'{}' successfully generated".format(pickup_file))
        
    return 0

        
if __name__ == "__main__":
    sys.exit(main())

