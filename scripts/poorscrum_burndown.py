#!/usr/bin/env python
# -*- coding: utf-8 -*-

__author__ = "sepp.heid@t-online.de"
__doc__ = """ """

import argparse
import os
import sys

from pptx import Presentation
from pptx.exc import PackageNotFoundError
from pptx.util import Inches

from poorscrum import SPRINT_FILE, BURNDOWN_SAVE_NAME, read_fields, __version__

import configparser

from matplotlib import rcParams
import matplotlib.pyplot as plt

import logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(os.path.basename(sys.argv[0]))



MIN_DAYS, MAX_DAYS = 0, 20
MIN_PERIODS, MAX_PERIODS = 0, 4
MIN_POINTS, MAX_POINTS = 0, 340 # 
          
def parse_arguments():
    """Parse command line arguments"""

    parser = argparse.ArgumentParser(
        prog = os.path.basename(sys.argv[0]),
        description = "Add burnout chart to the PPTX file",
        epilog = "Ensure the '.pptx/pptx_pickup.ini' file exists in the the home directory")
    
    parser.add_argument("--version", action = "version", version=__version__)

    parser.add_argument('-s', '--sprint_file', default=SPRINT_FILE,
                        help='SCRUM sprint setup file')

    parser.add_argument("-d", "--dry", required=False,
                        action="store_true", 
                        help="Do not create image, do not add to PPTX")

    parser.add_argument("the_pptx", nargs="?", 
                        help="The name of the PPTX file to create the burnout chart for")


    return parser.parse_args()


def parse_sprint(sprint_file = SPRINT_FILE):
    """ Read the scrum sprint file
    """
    sprint = configparser.ConfigParser()
    sprint.read(sprint_file)

    try:
        days = int(sprint.get("TIME", "days"))
    except configparser.NoSectionError:
        return None, None, None  # Illegal file content

    if days < MIN_DAYS or days > MAX_DAYS:
        return None, None # Illegal value

    try:
        periods = int(sprint.get("TIME", "periods"))
    except configparser.NoSectionError:
        return days, None, None  # Illegal file content

    if periods < MIN_PERIODS or periods > MAX_PERIODS:
        return days, None, None  # Illegal value

    try:
        points = 0
        for dev, point in sprint.items("TEAM"):
            if int(point) < MIN_POINTS or int(point) > MAX_POINTS:
                return days, periods, None 
            points += int(point)
    except:
        return days, periods, None # Illegal file content

    return days, periods, points


def extract_work(from_slide, with_pickup):
    """ Returns the the work to be done from the slide size items
    """

    text = ""
    for item, index in with_pickup.items():
        if (item != "size 1") and \
               (item != "size 2") and \
               (item != "size 3") and \
               (item != "size 4"):
            continue
            
        try:
            shape = from_slide.placeholders[int(index)]
        except KeyError:
            return list() # Does not contain story, no error
        
        if not shape.has_text_frame:
            continue

        paras = shape.text_frame.paragraphs
        if len(paras) > 1:
            return None # More than one paragraph, error
            
        for p in paras:
            for r in p.runs:
                for c in r.text.strip():
                    if not (c.isdigit() or c.isspace()):
                        return None #  Error
                    text += c
                text += ' '
    return [int(c) for c in text.split()]


def plot_burndown(points_left, last_edited, save_name = BURNDOWN_SAVE_NAME):
    sprint_days = len(points_left)

    rcParams['figure.figsize'] = (8, 6)

    plt.close("all")

    """ Draw the optimal line """
    plt.plot((0, sprint_days-1), (max(points_left), 0), color="blue", label="optimal")

    """ Draw the committed story points after sprint planning """
    plt.bar(0, points_left[0], color="green", label="start", width=0.5)

    """ Draw the actual work done """
    if (last_edited > 1) and (last_edited <= len(points_left)):
        plt.bar(range(1,last_edited), \
                points_left[1:last_edited], \
                color="red", label="actual", width=0.5)

    """ Draw the work still to be done """
    if (last_edited >= 1) and (last_edited < len(points_left)):
        plt.bar(range(last_edited,len(points_left)), \
                points_left[last_edited:len(points_left)], \
                color="grey", label="estimate", width=0.5)

    plt.ylim(0, points_left[0])
    plt.yticks(range(0, points_left[0]+1, max(1, int(points_left[0] / 5))))

    plt.xlim(-1, sprint_days)
    plt.xticks(range(0, sprint_days+1, max(5, int(sprint_days / 5))))

    plt.grid()
    plt.legend()

    plt.xlabel('SPRINT DAYS [days]')
    plt.ylabel('TO BE DONE [story points]')

    plt.title("Sprint Burndown", fontsize=20)
    plt.savefig(save_name)

    return save_name


def add_burndown_to_pptx(the_pptx, image_path):
    blank_slide_layout = the_pptx.slide_layouts[6]
    slide = the_pptx.slides.add_slide(blank_slide_layout)
    left = top = Inches(1)
    image = slide.shapes.add_picture(image_path, left, top)

def main():
    args = parse_arguments()

    if args.the_pptx is None:
        logger.error("Must provide PPTX backlog file name for work to be done!")
        return 1

    if args.sprint_file is None:
        logger.error("Must provide file with sprint setup!")
        return 2

    fields_map = read_fields()
    if fields_map is None:
        logger.error("Must provide correct story item mappings. Run 'pptx_learn.py'!")
        return 3

    days, periods, points = parse_sprint(args.sprint_file)
    if days is None:
        logger.error("Must provide scrum guide legal number of days for sprint!")
        return 4

    if periods is None:
        logger.error("Must provide scrum guide legal number of periods for sprint!")
        return 5

    if points is None:
        logger.error("Must provide scrum guide legal capacity of development team!")
        return 6

    logger.info("The sprint length is '{:d}' days.".format(days))
    logger.info("The sprint has '{:d}' periods with '{:d}' working days.".format(periods, int(days/periods)))
    logger.info("The team has a capacity '{:d}' story points.".format(points))

    try:
        prs = Presentation(args.the_pptx)
    except NameError:
        logger.error("Presentation has illegal name.")
        return 6
    except PackageNotFoundError:
        logger.error("Presentation is not found.")
        return 7
    
    logger.info("The Backlog '{}' has '{:d}' slides."
                .format(os.path.basename(args.the_pptx), len(prs.slides)))

    """
    The first day (day 0) in the sprint is the sprint planing day. All other
    days may be used fro working.

    The last value entered for a slide is used for the values not provided
    until the end of the sprint. These are called the 'unedited'
    """
    last_edited = days
    total_points_left = (days)*[0]
    for num, slide in enumerate(prs.slides):
        """ Calc work left for a slide in the sprint from sizes """
        slide_points_left = extract_work(slide, fields_map)
        if slide_points_left is None:
            logger.error("Slide '{:d}': Wrong syntax in size fields!".format(num+1))
            return 8

        if len(slide_points_left) == 0:
            logger.info("Slide '{:d}': Skipped since not a story file!".format(num+1))
            continue

        """ Determine the index of the last value entered """
        if len(slide_points_left) > last_edited:
            """ All sizes have to be equal for burndown """
            logger.error("Slide '{:d}': Size entries exceed previous ones.Fix sizes!".format(num+1))
            return 9

        last_edited = len(slide_points_left)
        if last_edited < 1:
            logger.error("Slide '{:d}': No Size entry.Fix sizes!".format(num+1))
            return 10

        """ Add entered values """
        edited_range = range(len(slide_points_left))
        for index in edited_range:
            total_points_left[index] += slide_points_left[index]

        """ Add the last entered value """
        unedited_range = range(len(slide_points_left), len(total_points_left))
        for index in unedited_range:
            total_points_left[index] += slide_points_left[-1]
        
        logger.info("Slide '{:d}': Included in work to be done!".format(num+1))

        
    logger.info("Work left is consistently entered including sprint day {:d}."
                .format(last_edited))

    """ Output work to be done in period format """
    nperiods = periods+1
    for period in range(int(days / nperiods)):
        period_estimate = total_points_left[(nperiods)*period:(nperiods)*(period+1)]
        logger.info("Work left estimate for period {:d}: {}".
                    format(period, str(period_estimate)))

    if points > total_points_left[0]:
        logger.info("Devs offer more points than required ('{:d}' > '{:d}'). Sprint can work!"
                    .format(points, total_points_left[0]))
        logger.info("'{:d}' points are available for analysis and spikes."
                    .format(points - total_points_left[0]))
    else:
        logger.warn("Devs offers less points than required ('{:d}' < '{:d}'). Sprint cannot work!"
                    .format(points, total_points_left[0]))
        logger.warn("Stories are to be reduced by '{:d}' points."
                    .format(total_points_left[0] - points))
    
    if args.dry:
        logger.info("Dry run finished: ok.")
        return 0

    save_name = plot_burndown(total_points_left, last_edited)
    logger.info("Saved the burddown chart to '{}'".format(save_name))

    add_burndown_to_pptx(prs, save_name)
    logger.info("Added the burddown chart to the presentation '{}'".format(args.the_pptx))

    try:
        prs.save(args.the_pptx)
    except:
        logger.error("Cannot save the presentation. If it is open consider to close it!")
        return 11
    
    logger.info("Saved the presentation to '{}'".format(args.the_pptx))
    
    return 0

        
if __name__ == "__main__":
    sys.exit(main())

